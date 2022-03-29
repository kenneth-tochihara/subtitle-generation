from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


with open('LB_script.txt') as f:
    script = f.readlines()

prs=Presentation("script-template.pptx")

for line in script:
    lyt=prs.slide_layouts[0] # choosing a slide layout
    slide=prs.slides.add_slide(lyt) # adding a slide
    title=slide.placeholders[0] # assigning a title
    title.text="".replace("\n", "")
    title_tf = title.text_frame
    p = title_tf.add_paragraph()
    
    # add a blank slide if indicated
    if "[BLANK]" in line:
        continue
    
    # see if there is a character role
    if line.count(":") == 0:
        p.text=line
        
    # if there isn't change the font
    else:
        run = p.add_run()
        run.text = line[:line.index(":")+1]
        font = run.font
        run.font.color.rgb = RGBColor(255, 51, 221)
        run = p.add_run()
        run.text = line[line.index(":")+1:]
    p.text.strip()


prs.save("script.pptx") # saving file

