from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

with open('script.txt') as f:
    script = f.readlines()

new_scripts = []
for index, line in enumerate(script):

    # remove parantheses if they exist
    new_line = line
    for idx in range(line.count("(")):
        print(index, line)
        extract = new_line[new_line.index("("):new_line.index(")")+1].strip()
        new_line = new_line.replace(extract, "")
    

    splits = []
    if new_line.count(":") == 2:

        temp = new_line[:new_line.rindex(':')]
        
        splitline_idx_list = []
        try: splitline_idx = temp.rindex('.')
        except ValueError:
            pass
        else: splitline_idx_list.append(splitline_idx)
        
        try: splitline_idx = temp.rindex('-')
        except ValueError:
            pass
        else: splitline_idx_list.append(splitline_idx)
        
        try: splitline_idx = temp.rindex('?')
        except ValueError:
            pass
        else: splitline_idx_list.append(splitline_idx)
        
        try: splitline_idx = temp.rindex('!')
        except ValueError:
            pass
        else: splitline_idx_list.append(splitline_idx)
        
        splitlines_left = temp[:splitline_idx+1]
        new_line = new_line[splitline_idx+1:]

        new_scripts.append(splitlines_left)
        new_scripts.append(new_line.strip())
        continue
    
    # remove if integer, null, or newline
    new_line = new_line.strip()
    if (new_line.strip().isnumeric()) or (new_line == ''):
        continue
    new_scripts.append(new_line.strip())

# for idx, line in enumerate(new_scripts):
#     if line.count(":") == 0:
#         print(line)
#     else:
#         print(line[:line.index(":")+1].strip(),"|", line[line.index(":")+1:].strip())

prs=Presentation("script.pptx")

for line in new_scripts:
    lyt=prs.slide_layouts[0] # choosing a slide layout
    slide=prs.slides.add_slide(lyt) # adding a slide
    title=slide.placeholders[0] # assigning a title
    title.text="".replace("\n", "")
    title_tf = title.text_frame
    p = title_tf.add_paragraph()
    
    # add a blank slide if indicated
    if line == "[BLANK]":
        continue
    
    # see if there is a character role
    if line.count(":") == 0:
        p.text=line
    # if there isn't change the font
    else:
        run = p.add_run()
        run.text = line[:line.index(":")+1]
        font = run.font
        run.font.color.rgb = RGBColor(238, 146, 143)

        run = p.add_run()
        run.text = line[line.index(":")+1:]
    p.text.strip()


prs.save("script.pptx") # saving file

# with open('new_script.txt', 'w') as f:
#     f.writelines(new_scripts)